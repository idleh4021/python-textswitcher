import os
from pptx import Presentation
from pathlib import Path
from tkinter import filedialog
import sys

def replace_text_in_ppt(input_path, output_path, old_text, new_text):
    prs = Presentation(input_path)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)
    
    prs.save(output_path)
    print(f"파일 저장 완료: {output_path}")

def replace_text(prs, old_text, new_text):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)
                    else:
                        # old_text가 run에 안 걸쳐있으면 넘어감
                        pass

def print_ppt_text_structure(prs):
    for slide_idx, slide in enumerate(prs.slides, start=1):
        print(f"\n--- 슬라이드 {slide_idx} ---")
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if not shape.has_text_frame:
                continue
            print(f"  도형 {shape_idx} 텍스트 프레임:")
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                para_text = "".join(run.text for run in paragraph.runs)
                print(f"    문단 {para_idx}: '{para_text}'")
                for run_idx, run in enumerate(paragraph.runs, start=1):
                    print(f"      런 {run_idx}: '{run.text}'")

def replace_text_across_runs(prs, old_text, new_text):
    for slide in prs.slides:
        for shape in slide.shapes:
            # ✅ 테이블이면 테이블 셀 텍스트도 순회
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        full_text = ""
                        runs = []
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                runs.append(run)
                                full_text += run.text

                        if old_text in full_text:
                            full_text = full_text.replace(old_text, new_text)
                            for run in runs:
                                run.text = ""
                            if runs:
                                runs[0].text = full_text

            # ✅ 일반 텍스트 프레임 도형
            elif shape.has_text_frame:
                text_frame = shape.text_frame
                full_text = ""
                runs = []
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        runs.append(run)
                        full_text += run.text
                #print(f'일반텍스트(bool, old,new):in {old_text in full_text} , {old_text}, {full_text}')
                if old_text in full_text:
                    
                    full_text = full_text.replace(old_text, new_text)
                    for run in runs:
                        run.text = ""
                    if runs:
                        runs[0].text = full_text

                    
def replace_text_across_runs_by_dict(prs, switch_words):
    for i in switch_words:
        old_text = i
        new_text = switch_words[i]
        replace_text_across_runs(prs,old_text,new_text)


if __name__ == "__main__":
    # 예: C 드라이브 내 'sample.pptx' 파일 불러오기
    
    load_files = filedialog.askopenfilenames(filetypes=[('ppt Files','*.ppt;*.pptx')])
    switch_words={}
    print('치환할 문자열을 입력해주세요 (예시: 코끼리 -> 냉장고 치환 : 코끼리 냉장고) \n치환할 문자가 더이상 없으면 엔터를 입력하세요 ')
    while True:
        print('치환 문자:',end='',flush=True)
        s = sys.stdin.readline().strip()
        print(f's={s}')
        if s=='':break
        else:
            s,t = s.split()
            switch_words[s]=t
    print(switch_words)
    #exit()
    #print(load_files)
    for load_file in load_files :
        
        # 내문서 내 TextSwitcher 폴더 경로
        documents_path = os.path.join(os.path.expanduser("~"), "Documents")
        output_folder = os.path.join(documents_path, "TextSwitcher")
        os.makedirs(output_folder, exist_ok=True)

        output_ppt_path = os.path.join(output_folder, os.path.basename(load_file))

        prs = Presentation(load_file)
        print_ppt_text_structure(prs)
        replace_text_across_runs_by_dict(prs,switch_words)
        #replace_text(prs,old_text_to_replace,new_text_to_replace)
        try:
            prs.save(output_ppt_path)
            print('저장이 완료되었습니다.')
            os.startfile(output_folder)
        except Exception as e:
            print('저장에 실패했습니다.파일이 열려있는지 확인하세요')
        
        #replace_text_in_ppt(input_ppt_path, output_ppt_path, old_text_to_replace, new_text_to_replace)
